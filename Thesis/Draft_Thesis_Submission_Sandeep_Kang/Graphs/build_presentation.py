from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ---------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------
TITLE = "Machine Learning-Based Forecasting of Ride-Hailing Demand, Revenue, and Operations: A Case Study of Delhi NCR, India"
SUBTITLE = "INFO-I 492 Senior Thesis — Indiana University"
AUTHOR = "Sandeep Kang"

# Put your figures in ./Graphs/ with these filenames
FIG_DIR = "Graphs"
FIG_MAP = {
    # Intro and Background
    3: "Figure1_Problem_Context_Map.png",                    # Slide 3 (optional)
    4: "Figure2_Project_Workflow_Overview.png",              # Slide 4

    # Literature Review
    6: "Figure3_ModelPerformanceComparison.png",             # Slide 6
    7: "Figure5_Fairness_Sustainability_Governance.png",     # Slide 7 (conceptual)

    # Methods
    9:  "Figure7_Data_Preprocessing_Workflow.png",           # Slide 9
    10: "Figure8_Feature_Engineering_Pipeline.png",          # Slide 10
    11: "Figure9_Model_Training_Evaluation_Pipeline.png",    # Slide 11
    12: "Figure4_Dynamic_Pricing_System_Architecture.png",   # Slide 12
    13: "Figure6_Conceptual_Framework.png",                  # Slide 13

    # Results
    15: "Figure11_Hourly_Weekly_Demand_Heatmap.png",         # Slide 15
    16: "Figure10_Combined_Performance_And_Features.png",    # Slide 16
    17: "Figure12_Fare_Prediction_Comparison.png",           # Slide 17
    18: "Figure14_Demand_Fulfillment_Cancellation.png",      # Slide 18
    19: "Figure15_Regional_Ride_Flow.png",                   # Slide 19
    20: "Figure13_Surge_Distribution_Heatmap.png",           # Slide 20
    21: "Figure16_Emission_Reduction.png",                   # Slide 21
}

OUTPUT_NAME = "Machine_Learning_Based_Forecasting_RideHailing_DelhiNCR.pptx"

# ---------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------
def set_title_style(text_frame, font_size=36, bold=True, color=(20, 33, 61)):
    p = text_frame.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    p.alignment = PP_ALIGN.LEFT

def set_body_style(text_frame, font_size=20, color=(33, 37, 41)):
    for p in text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = Pt(font_size)
            r.font.color.rgb = RGBColor(*color)

def add_image_if_exists(slide, image_path, left=Inches(0.5), top=Inches(1.8), width=Inches(9.0)):
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=width)
            return True
        except Exception:
            return False
    return False

def add_simple_bullets(shape, bullets):
    tf = shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
    set_body_style(tf)

def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

def add_section_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    title.text = title_text
    set_title_style(title.text_frame)
    return slide

def add_title_and_content(prs, title_text, body_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = title_text
    set_title_style(title.text_frame)
    body = slide.placeholders[1]
    if body_lines:
        add_simple_bullets(body, body_lines)
    else:
        add_simple_bullets(body, ["Add main points here..."])
    return slide

# ---------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = TITLE
set_title_style(slide.shapes.title.text_frame, font_size=34)
subtitle = slide.placeholders[1]
subtitle.text = f"{AUTHOR}\n{SUBTITLE}"
set_body_style(subtitle.text_frame, font_size=18)
add_notes(slide, "Introduce yourself and the thesis title. Provide a quick overview of the agenda and duration.")

# 2. Agenda
s2 = add_title_and_content(prs, "Agenda", [
    "Background and problem context",
    "Hypothesis and research questions",
    "Literature review highlights",
    "Methods and dataset",
    "Results",
    "Discussion and implications",
    "Conclusion and future work"
])
add_notes(s2, "State the flow. Set expectations for figures and transitions. Target one to two minutes.")

# 3. Background - Problem Context (Figure 1 optional)
s3 = add_title_and_content(prs, "Background: Ride-Hailing in Delhi NCR", [
    "Scale of demand and service variability",
    "Operational challenges: cancellations and fulfillment",
    "Motivation for machine learning-based forecasting"
])
img3 = os.path.join(FIG_DIR, FIG_MAP.get(3, ""))
add_image_if_exists(s3, img3)
add_notes(s3, "Use the Problem Context Map if available. Ground the audience in the local context and stakes.")

# 4. Project Workflow Overview (Figure 2)
s4 = add_title_and_content(prs, "Project Workflow Overview (Figure 2)", [
    "Data intake and cleaning",
    "Feature engineering and model training",
    "Evaluation and operational integration"
])
img4 = os.path.join(FIG_DIR, FIG_MAP.get(4, ""))
add_image_if_exists(s4, img4)
add_notes(s4, "Walk through the end-to-end pipeline. Keep it high-level here; details will come later.")

# 5. Hypothesis and Research Questions
s5 = add_title_and_content(prs, "Hypothesis and Research Questions", [
    "Hypothesis: Integrated ML pipeline improves predictive accuracy and operational insight.",
    "RQ1: Can we forecast demand and revenue accurately at hourly scale?",
    "RQ2: Can forecasts inform pricing and resource allocation decisions?"
])
add_notes(s5, "Make the claims precise and testable. Keep this slide crisp.")

# 6. Literature Review: Forecasting Models (Figure 3)
s6 = add_title_and_content(prs, "Literature Review: Forecasting Models", [
    "Classical vs deep learning performance differences",
    "Spatiotemporal and attention-based advances",
    "Evidence for lower MAE and RMSE using deep and ensemble models"
])
img6 = os.path.join(FIG_DIR, FIG_MAP.get(6, ""))
add_image_if_exists(s6, img6)
add_notes(s6, "Point to how your results align or differ. Mention key citations verbally as needed.")

# 7. Literature Review: Fairness and Governance (Figure 5)
s7 = add_title_and_content(prs, "Literature Review: Fairness and Governance", [
    "Fairness-aware optimization and constraints",
    "Regulatory expectations and auditability",
    "Implications for pricing and accessibility"
])
img7 = os.path.join(FIG_DIR, FIG_MAP.get(7, ""))
add_image_if_exists(s7, img7)
add_notes(s7, "Use the conceptual diagram to explain the intersection of fairness, sustainability, and policy.")

# 8. Methods Overview
s8 = add_title_and_content(prs, "Methods Overview", [
    "Dataset: trip-level records with time, locations, vehicle type, fares",
    "Tools: pandas, scikit-learn, gradient boosting, random forests",
    "Evaluation: MAE, RMSE, R-squared; visual diagnostics"
])
add_notes(s8, "Preview the steps; specifics follow on subsequent slides.")

# 9. Data Preprocessing Workflow (Figure 7)
s9 = add_title_and_content(prs, "Data Preprocessing Workflow (Figure 7)", [
    "Cleaning and type coercion",
    "Handling missing values and duplicates",
    "Output: cleaned dataset for feature generation"
])
img9 = os.path.join(FIG_DIR, FIG_MAP.get(9, ""))
add_image_if_exists(s9, img9)
add_notes(s9, "Call out any caveats or filtering choices that affect downstream results.")

# 10. Feature Engineering Pipeline (Figure 8)
s10 = add_title_and_content(prs, "Feature Engineering Pipeline (Figure 8)", [
    "Fare per km, tip ratio, peak hour indicators",
    "Encodings for payment and day-of-week",
    "Final training matrix"
])
img10 = os.path.join(FIG_DIR, FIG_MAP.get(10, ""))
add_image_if_exists(s10, img10)
add_notes(s10, "Mention rationale for each engineered feature and any ablations you considered.")

# 11. Model Training and Evaluation (Figure 9)
s11 = add_title_and_content(prs, "Model Training and Evaluation (Figure 9)", [
    "Baselines and advanced models",
    "Train-test splits and validation",
    "Performance metrics and logging"
])
img11 = os.path.join(FIG_DIR, FIG_MAP.get(11, ""))
add_image_if_exists(s11, img11)
add_notes(s11, "Emphasize reproducibility. Note any hyperparameters and selection criteria.")

# 12. Dynamic Pricing System Architecture (Figure 4)
s12 = add_title_and_content(prs, "Dynamic Pricing System Architecture (Figure 4)", [
    "Demand prediction and driver availability",
    "Pricing feedback loops and constraints",
    "Operational alignment"
])
img12 = os.path.join(FIG_DIR, FIG_MAP.get(12, ""))
add_image_if_exists(s12, img12)
add_notes(s12, "Show how predictions could inform price multipliers. Keep it conceptual, not proprietary.")

# 13. Conceptual Framework (Figure 6)
s13 = add_title_and_content(prs, "Conceptual Framework (Figure 6)", [
    "Integrated forecasting, revenue, and operations",
    "Data and feedback flows",
    "External factors influencing modules"
])
img13 = os.path.join(FIG_DIR, FIG_MAP.get(13, ""))
add_image_if_exists(s13, img13)
add_notes(s13, "Explain bidirectional flows and one-way external inputs clearly.")

# 14. Results Overview
s14 = add_title_and_content(prs, "Results Overview", [
    "Temporal demand patterns",
    "Model accuracy comparisons",
    "Operational and pricing insights"
])
add_notes(s14, "Set up the next sequence of detailed results slides.")

# 15. Results: Temporal Patterns (Figure 11)
s15 = add_title_and_content(prs, "Results: Hourly and Weekly Demand (Figure 11)", [
    "Weekday peak patterns vs weekend distribution",
    "Implications for staffing and availability",
    "Variability and potential drivers"
])
img15 = os.path.join(FIG_DIR, FIG_MAP.get(15, ""))
add_image_if_exists(s15, img15)
add_notes(s15, "Call out the specific peaks. Tie to actionable decisions.")

# 16. Results: Combined Performance and Features (Figure 10)
s16 = add_title_and_content(prs, "Results: Model Performance and Features (Figure 10)", [
    "Comparative accuracy across models",
    "Feature importances and SHAP insights",
    "Interpretability and stability"
])
img16 = os.path.join(FIG_DIR, FIG_MAP.get(16, ""))
add_image_if_exists(s16, img16)
add_notes(s16, "Discuss why some features dominate and how that affects policy and operations.")

# 17. Results: Fare Prediction (Figure 12)
s17 = add_title_and_content(prs, "Results: Fare Prediction Comparison (Figure 12)", [
    "Error distribution across models",
    "Robustness to demand spikes",
    "Generalization considerations"
])
img17 = os.path.join(FIG_DIR, FIG_MAP.get(17, ""))
add_image_if_exists(s17, img17)
add_notes(s17, "Note any heteroskedasticity or tail risks apparent in the comparisons.")

# 18. Results: Fulfillment and Cancellations (Figure 14)
s18 = add_title_and_content(prs, "Results: Fulfillment and Cancellations (Figure 14)", [
    "Service reliability by hour",
    "Driver and rider cancellation patterns",
    "Targets for mitigation"
])
img18 = os.path.join(FIG_DIR, FIG_MAP.get(18, ""))
add_image_if_exists(s18, img18)
add_notes(s18, "If the expected dips were not visible, explain adjustments or alternate metrics used.")

# 19. Results: Regional Ride Flow (Figure 15)
s19 = add_title_and_content(prs, "Results: Regional Ride Flow (Figure 15)", [
    "Flows between central and peripheral zones",
    "Implications for repositioning and supply",
    "Equity and accessibility considerations"
])
img19 = os.path.join(FIG_DIR, FIG_MAP.get(19, ""))
add_image_if_exists(s19, img19)
add_notes(s19, "Explain how zones were defined and how to interpret the heat intensity.")

# 20. Results: Surge and Pricing Patterns (Figure 13)
s20 = add_title_and_content(prs, "Results: Surge and Pricing Patterns (Figure 13)", [
    "Surge distribution across time and zones",
    "Elasticity and user abandonment thresholds",
    "Policy guardrails"
])
img20 = os.path.join(FIG_DIR, FIG_MAP.get(20, ""))
add_image_if_exists(s20, img20)
add_notes(s20, "Relate back to literature on elasticity and perceived fairness.")

# 21. Results: Emission Reduction Scenario (Figure 16)
s21 = add_title_and_content(prs, "Results: Emission Reduction Scenario (Figure 16)", [
    "Idle time and pooling impacts",
    "EV scheduling implications",
    "Co-benefits and trade-offs"
])
img21 = os.path.join(FIG_DIR, FIG_MAP.get(21, ""))
add_image_if_exists(s21, img21)
add_notes(s21, "State assumptions and whether the scenario is illustrative or empirically derived.")

# 22. Discussion
s22 = add_title_and_content(prs, "Discussion", [
    "How results support or qualify the hypothesis",
    "Operational and policy implications",
    "Limitations and interpretability considerations"
])
add_notes(s22, "Use this slide to weave results back to your core research questions.")

# 23. Conclusion
s23 = add_title_and_content(prs, "Conclusion", [
    "Integrated ML pipeline improved accuracy and decision linkage",
    "Context matters for transferability",
    "Balance between accuracy and transparency"
])
add_notes(s23, "Keep concise and memorable. This sets up future work.")

# 24. Future Work
s24 = add_title_and_content(prs, "Future Work", [
    "Regional generalization and transfer learning",
    "Real-time closed-loop optimization",
    "Fairness audits and explainable interfaces"
])
add_notes(s24, "Tie future work to real constraints you observed.")

# 25. Q and A
s25 = add_title_and_content(prs, "Questions and Answers", [
    "Methods, data, ethics, operationalization",
    "Trade-offs and deployment considerations",
    "Thank you"
])
add_notes(s25, "Invite questions. Keep backup slides handy if you add them later.")

# Save
prs.save(OUTPUT_NAME)
print(f"Saved: {OUTPUT_NAME}")